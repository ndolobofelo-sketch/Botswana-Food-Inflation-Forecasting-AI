"""
Food Inflation Forecasting
Final Hackathon Presentation Generator
"""

from pptx import Presentation
from pptx.util import Inches, Pt
import os


os.makedirs(
    "reports",
    exist_ok=True
)


prs = Presentation()


def add_slide(title, bullets):

    slide = prs.slides.add_slide(
        prs.slide_layouts[1]
    )

    slide.shapes.title.text = title

    body = slide.placeholders[1]

    body.text = bullets


    for paragraph in body.text_frame.paragraphs:
        paragraph.font.size = Pt(18)



# --------------------------------------------------
# Slide 1
# --------------------------------------------------

add_slide(
    "VenturePulse: Food Price Inflation Forecasting in Botswana",

"""
AI-powered monthly food price forecasting system

Objective:
Predict future food price movements using
multi-source economic and regional data.

Approach:
Data Engineering + Machine Learning + Explainable AI
"""
)



# --------------------------------------------------
# Slide 2
# --------------------------------------------------

add_slide(
    "Problem & Data Sources",

"""
Problem:
Food price volatility creates challenges for
households, businesses and policymakers.

Integrated datasets:

• FAO Botswana food prices
• Baltic Dry Index shipping data
• Brent crude oil prices
• Botswana policy rate
• Human Capital Project indicators

Period:
January 2000 - December 2023
"""
)



# --------------------------------------------------
# Slide 3
# --------------------------------------------------

add_slide(
    "Data Pipeline & Feature Engineering",

"""
Pipeline:

Raw Data
    ↓
Cleaning & Validation
    ↓
Monthly Alignment
    ↓
Feature Engineering
    ↓
Machine Learning Model


Created features:

• Lag variables
• Rolling averages
• Economic changes
• Seasonal time features
• Regional food indicators
"""
)



# --------------------------------------------------
# Slide 4
# --------------------------------------------------

add_slide(
    "Machine Learning Results",

"""
Models evaluated:

1. Linear Regression
2. Random Forest
3. Gradient Boosting


Selected Model:

Gradient Boosting Regressor


Reason:

Best forecasting performance and ability
to capture nonlinear relationships.
"""
)



# --------------------------------------------------
# Slide 5
# --------------------------------------------------

add_slide(
    "2024 Forecast & Impact",

"""
Generated forecast:

January 2024 - December 2024


Expected impact:

• Support food security planning
• Improve economic forecasting
• Assist policy decisions
• Demonstrate AI-driven analytics
"""
)



output = (
    "reports/"
    "Food_Inflation_Forecasting_Presentation.pptx"
)


prs.save(output)


print("="*80)
print("PRESENTATION CREATED")
print("="*80)

print(output)