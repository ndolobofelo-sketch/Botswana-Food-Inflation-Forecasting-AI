from pptx import Presentation
from pptx.util import Inches, Pt
import os


print("=" * 70)
print("CREATING FINAL FOOD INFLATION PRESENTATION")
print("=" * 70)


output = "reports/presentation/Food_Inflation_Final_Presentation.pptx"

os.makedirs("reports/presentation", exist_ok=True)


prs = Presentation()


def add_slide(title, bullets):

    slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(slide_layout)

    slide.shapes.title.text = title

    body = slide.placeholders[1]

    body.text = bullets


    for paragraph in body.text_frame.paragraphs:
        paragraph.font.size = Pt(20)


# Slide 1

add_slide(
    "Forecasting Botswana Food Inflation Using Machine Learning",

    """
Problem:
• Food price inflation creates economic pressure on households and policymakers.

Objective:
• Develop an AI forecasting framework to predict food inflation trends.
• Identify key economic drivers affecting food prices.
• Support evidence-based decision making.
"""
)


# Slide 2

add_slide(
    "Data Sources and Integration",

    """
Integrated datasets:

• FAO Botswana food price indicators
• Baltic Dry Index (global shipping costs)
• Brent crude oil prices
• Botswana policy rate
• Human Capital Project indicators

Data preparation:
• Time alignment
• Monthly aggregation
• Dataset merging
• Missing value handling
"""
)


# Slide 3

add_slide(
    "Feature Engineering Pipeline",

    """
Created predictive features:

• Historical lag variables
• Rolling averages
• Time-based features
• Economic indicator relationships

Pipeline:

Raw Data
    ↓
Cleaning & Transformation
    ↓
Feature Creation
    ↓
Machine Learning Models
    ↓
Forecast Generation
"""
)


# Slide 4

add_slide(
    "Model Performance and Interpretability",

    """
Models evaluated:

• XGBoost
• Random Forest
• LSTM Neural Network

Interpretability:

• Feature importance analysis
• Economic driver identification
• Regression linkage analysis

Evidence produced:
• Model comparison
• Coefficients
• P-values
• R-squared statistics
"""
)


# Slide 5

add_slide(
    "Forecast Results and Conclusion",

    """
Forecast outcome:

• Machine learning models provide forward food inflation predictions.
• Historical economic relationships improve forecast interpretation.

Key value:

• Supports policymakers
• Improves inflation monitoring
• Provides data-driven insights for Botswana food security

Conclusion:

AI forecasting combined with economic analysis creates a stronger
framework for understanding future food price movements.
"""
)


prs.save(output)


print("=" * 70)
print("PRESENTATION CREATED SUCCESSFULLY")
print(output)
print("=" * 70)